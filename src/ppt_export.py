# src/ppt_export.py
from pptx import Presentation
from pptx.util import Inches

def create_ppt(summary_dfs: dict):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Monthly Board KPI Report"
    slide.placeholders[1].text = "Generated automatically"

    for title, df in summary_dfs.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        rows, cols = df.shape[0]+1, df.shape[1]
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        for col_idx, col_name in enumerate(df.columns):
            table.cell(0, col_idx).text = col_name

        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx+1, col_idx).text = str(value)
    return prs
