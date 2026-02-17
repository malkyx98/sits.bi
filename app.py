import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

# -----------------------------
# PAGE CONFIG
# -----------------------------
st.set_page_config(page_title="AI-Powered Excel Analyzer", layout="wide")
st.title("🤖 AI-Powered Global Excel Analyzer")
st.write("Upload ANY Excel file → Get instant AI-driven KPIs, trends, and reports")

# -----------------------------
# FILE UPLOAD
# -----------------------------
uploaded_file = st.file_uploader("Upload Excel File", type=["xlsx","xls"])
if uploaded_file is None:
    st.info("Please upload an Excel file to start analysis.")
    st.stop()

# Read Excel safely
try:
    df = pd.read_excel(uploaded_file)
except Exception as e:
    st.error(f"Failed to read Excel file: {e}")
    st.stop()

st.success(f"File uploaded! Rows: {df.shape[0]}, Columns: {df.shape[1]}")

# Clean column names
df.columns = [str(c).strip() if c else f"Column_{i}" for i,c in enumerate(df.columns)]

with st.expander("Preview Data"):
    st.dataframe(df.head())

# -----------------------------
# AI COLUMN ROLE DETECTION
# -----------------------------
def detect_column_roles(df):
    numeric_cols, categorical_cols, date_cols, id_cols, status_cols = [], [], [], [], []

    for col in df.columns:
        sample = df[col].dropna()
        if sample.empty:
            continue

        # Date detection
        try:
            parsed = pd.to_datetime(sample, errors='coerce')
            if parsed.notna().sum() / len(sample) > 0.5:
                date_cols.append(col)
                continue
        except:
            pass

        # Numeric detection
        numeric_series = pd.to_numeric(sample, errors='coerce')
        if numeric_series.notna().any():
            numeric_cols.append(col)
            continue

        # Status detection (yes/no, closed/open, passed/failed)
        lower = sample.astype(str).str.lower()
        if lower.isin(['yes','no','closed','open','passed','failed']).any():
            status_cols.append(col)
            continue

        # ID / identifier detection (names, codes)
        if any(keyword in col.lower() for keyword in ['id','name','ref','agent','technician','caller']):
            id_cols.append(col)
            continue

        # Default categorical
        categorical_cols.append(col)

    return numeric_cols, categorical_cols, date_cols, id_cols, status_cols

numeric_cols, categorical_cols, date_cols, id_cols, status_cols = detect_column_roles(df)

st.subheader("Detected Column Roles")
c1,c2,c3,c4,c5 = st.columns(5)
c1.info(f"🔢 Numeric: {numeric_cols if numeric_cols else 'None'}")
c2.warning(f"🏷️ Categorical: {categorical_cols if categorical_cols else 'None'}")
c3.success(f"📅 Date: {date_cols if date_cols else 'None'}")
c4.primary(f"🆔 ID/Name: {id_cols if id_cols else 'None'}")
c5.secondary(f"✅ Status: {status_cols if status_cols else 'None'}")

# -----------------------------
# AI AUTOMATIC ANALYSIS
# -----------------------------
st.subheader("📊 Automatic KPI & Summary Tables")

# Numeric KPIs
if numeric_cols:
    for col in numeric_cols:
        series = pd.to_numeric(df[col], errors='coerce')
        k1,k2,k3,k4 = st.columns(4)
        k1.metric(f"{col} Total", round(series.sum(),2))
        k2.metric(f"{col} Avg", round(series.mean(),2))
        k3.metric(f"{col} Max", round(series.max(),2))
        k4.metric(f"{col} Min", round(series.min(),2))
else:
    st.info("No numeric columns detected.")

# Status KPIs
if status_cols:
    for col in status_cols:
        counts = df[col].value_counts()
        st.markdown(f"**{col} Summary**")
        st.table(counts)

# Categorical top 10
if categorical_cols:
    for col in categorical_cols:
        counts = df[col].astype(str).value_counts().head(10)
        if not counts.empty:
            st.markdown(f"**{col} Top 10 Values**")
            st.bar_chart(counts)

# Date trends
if date_cols:
    for col in date_cols:
        try:
            df[col+'_parsed'] = pd.to_datetime(df[col], errors='coerce')
            trend = df.groupby(df[col+'_parsed'].dt.to_period("M")).size().reset_index(name='Count')
            trend[col+'_parsed'] = trend[col+'_parsed'].astype(str)
            fig = px.line(trend, x=col+'_parsed', y='Count', title=f"Trend over time ({col})")
            st.plotly_chart(fig, use_container_width=True)
        except:
            continue

# -----------------------------
# CORRELATION FOR NUMERIC
# -----------------------------
if len(numeric_cols) > 1:
    st.subheader("🔥 Numeric Correlation Heatmap")
    corr = df[numeric_cols].apply(pd.to_numeric, errors='coerce').corr()
    fig_corr = px.imshow(corr, text_auto=True, color_continuous_scale='RdBu_r', title="Correlation Heatmap")
    st.plotly_chart(fig_corr, use_container_width=True)

# -----------------------------
# DOWNLOAD PROCESSED EXCEL & PPT
# -----------------------------
st.subheader("📥 Download Processed Reports")

# Excel
excel_buffer = BytesIO()
with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
    df.to_excel(writer, sheet_name='Processed_Data', index=False)
excel_buffer.seek(0)
st.download_button("Download Excel Report", excel_buffer, "AI_Excel_Report.xlsx")

# PowerPoint
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "AI-Powered Excel Report"
rows, cols_table = df.shape
table = slide.shapes.add_table(min(rows+1, 20), min(cols_table, 10), Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
for j, col_name in enumerate(df.columns[:10]):
    table.cell(0,j).text = str(col_name)
for i, row in enumerate(df.head(20).values):
    for j, val in enumerate(row[:10]):
        table.cell(i+1,j).text = str(val)
ppt_buffer = BytesIO()
prs.save(ppt_buffer)
ppt_buffer.seek(0)
st.download_button("Download PowerPoint Preview", ppt_buffer, "AI_Excel_Report.pptx")

